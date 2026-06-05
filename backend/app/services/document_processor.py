import os
import shutil
from pathlib import Path

def get_file_extension(filename: str) -> str:
    return Path(filename).suffix.lower()

def extract_text_from_pdf(file_path: str) -> dict:
    """Extract text from PDF using pdfplumber. If a page has no text (scanned), note it."""
    import pdfplumber
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            pages.append({"page": i + 1, "text": text.strip()})
    
    full_text = "\n\n".join(p["text"] for p in pages if p["text"])
    return {
        "text": full_text,
        "pages": len(pages),
        "file_type": "pdf"
    }

def extract_text_from_docx(file_path: str) -> dict:
    """Extract text from DOCX using python-docx."""
    from docx import Document
    doc = Document(file_path)
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())
    
    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    
    return {
        "text": "\n".join(paragraphs),
        "pages": 1,
        "file_type": "docx"
    }

def extract_text_from_pptx(file_path: str) -> dict:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
    
    return {
        "text": "\n\n".join(slides_text),
        "pages": len(prs.slides),
        "file_type": "pptx"
    }

def process_document(file_path: str, filename: str) -> dict:
    """Process document based on file extension."""
    ext = get_file_extension(filename)
    
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in (".docx", ".doc"):
        return extract_text_from_docx(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path)
    elif ext in (".jpg", ".jpeg", ".png", ".bmp", ".tiff"):
        return {"text": None, "pages": 1, "file_type": "image"}  # Will be handled by OCR
    else:
        raise ValueError(f"Unsupported file format: {ext}")
